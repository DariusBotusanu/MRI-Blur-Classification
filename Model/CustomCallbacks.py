import os
import keras
import pandas as pd
import seaborn as sns

import pptx
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE

sns.set_style('darkgrid')
sns.set(rc={'figure.figsize':(8,6)})

class PlottingCallback(keras.callbacks.Callback):

    def __init__(self, folder_path):
      self.folder_path = folder_path
      super(PlottingCallback, self).__init__()

    def on_train_begin(self, logs=None):
        pass

    def on_train_end(self, logs=None):
        files = os.listdir(self.folder_path)
        if 'training.log' not in files:
          print("WARNING: There is no training.log in the specified folder")
          return
        
        df = pd.read_csv(self.folder_path+'/training.log')
        df.set_index('epoch', inplace=True)
        loss_plot = df[['loss','val_loss']].plot()
        acc_plot = df[['accuracy', 'val_accuracy']].plot()

        loss_plot_fig = loss_plot.get_figure()
        acc_plot_fig = acc_plot.get_figure()

        loss_plot_fig.savefig(self.folder_path+'/loss.jpg')
        acc_plot_fig.savefig(self.folder_path+'/accuracy.jpg')
        

    def on_epoch_begin(self, epoch, logs=None):
        pass

    def on_epoch_end(self, epoch, logs=None):
        pass

    def on_test_begin(self, logs=None):
        pass

    def on_test_end(self, logs=None):
        pass

    def on_predict_begin(self, logs=None):
        pass

    def on_predict_end(self, logs=None):
        pass

    def on_train_batch_begin(self, batch, logs=None):
        pass

    def on_train_batch_end(self, batch, logs=None):
        pass

    def on_test_batch_begin(self, batch, logs=None):
        pass

    def on_test_batch_end(self, batch, logs=None):
        pass

    def on_predict_batch_begin(self, batch, logs=None):
        pass

    def on_predict_batch_end(self, batch, logs=None):
        pass

class PptxReportCallback(keras.callbacks.Callback):

    def __init__(self, folder_path, model_dir):
      self.folder_path = folder_path
      self.model_dir = model_dir
      super(PptxReportCallback, self).__init__()

    def on_train_begin(self, logs=None):
        pass

    def on_train_end(self, logs=None):
        files = os.listdir(self.folder_path)
        if 'model_report.pptx' not in files:
          print("WARNING: No pptx file was found!")
          return
        
        prs = Presentation(self.folder_path+'/model_report.pptx')

        #We add a new blank slide
        SLD_LAYOUT_TITLE_AND_CONTENT = 6
        slide_layout = prs.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
        slide = prs.slides.add_slide(slide_layout) #ADDS a new blank slide

        #SHAPES OF THE ADDED SLIDE
        shapes = prs.slides[len(prs.slides)-1].shapes 

        #We add the loss and accuracy images to the slide
        loss = shapes.add_picture(self.model_dir+'/loss.jpg',height=Inches(3.79),width=Inches(5.05),left=0,top=0)
        acc = shapes.add_picture(self.model_dir+'/accuracy.jpg',height=Inches(3.79),width=Inches(5.05),left=0,top=Inches(3.79))

        #We add the model summary to the slide
        summary = shapes.add_textbox(left=Inches(5.05), height=Inches(7.5), width=Inches(8.28),top=0)

        #We prepare the 'summary' shape for writing
        summary_text_frame = summary.text_frame
        summary_text_frame.word_wrap = True
        summary_paragraph = summary_text_frame.add_paragraph()

        #We read the optimizer
        optimizer_file = open(self.model_dir+'/optimizer.txt', 'r')
        optimizer_text = optimizer_file.read()
        optimizer_file.close()

        #We read thhe model summary
        summary_file = open(self.model_dir+'/model_summary.txt', 'r')
        summary_text = summary_file.read()
        summary_file.close()
        lines_list = summary_text.split('\n')

        #aesthetics
        for i in range(len(lines_list)):
          if '==' in lines_list[i] or '__' in lines_list[i]:
            lines_list[i] = lines_list[i][:len(lines_list[i])//2]

        lines_list[0] = 'Model: '+self.model_dir[len('Callbacks/'):] #we rename the model


        summary_text = ''
        for line in lines_list:
          summary_text += line + '\n' 
        
        summary_text += 'Optimizer: '+optimizer_text
        
        #We write the actual summary to the shape
        summary_paragraph.text = summary_text

        #We fit the text in the shape (it will need refreshing in the pptx)
        summary_paragraph.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        summary.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        #We overwrite the pptx file
        prs.save(self.folder_path+'/model_report.pptx')

    def on_epoch_begin(self, epoch, logs=None):
        pass

    def on_epoch_end(self, epoch, logs=None):
        pass

    def on_test_begin(self, logs=None):
        pass

    def on_test_end(self, logs=None):
        pass

    def on_predict_begin(self, logs=None):
        pass

    def on_predict_end(self, logs=None):
        pass

    def on_train_batch_begin(self, batch, logs=None):
        pass

    def on_train_batch_end(self, batch, logs=None):
        pass

    def on_test_batch_begin(self, batch, logs=None):
        pass

    def on_test_batch_end(self, batch, logs=None):
        pass

    def on_predict_batch_begin(self, batch, logs=None):
        pass

    def on_predict_batch_end(self, batch, logs=None):
        pass